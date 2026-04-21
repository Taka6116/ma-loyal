"""One-off: extract text from PPTX/PDF to Markdown for archival."""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pypdf import PdfReader


def pptx_to_md(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = [f"# {path.stem}\n"]
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n## スライド {i}\n\n")
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if texts:
            parts.append("\n\n".join(texts))
            parts.append("\n")
        else:
            parts.append("（テキストなし）\n")
    return "".join(parts).strip() + "\n"


def pdf_to_md(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: list[str] = [f"# {path.stem}\n"]
    for i, page in enumerate(reader.pages, start=1):
        parts.append(f"\n## ページ {i}\n\n")
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            parts.append(text)
            parts.append("\n")
        else:
            parts.append("（抽出テキストなし — 画像主体の可能性）\n")
    return "".join(parts).strip() + "\n"


def sanitize_filename(stem: str) -> str:
    bad = r'[<>:"/\\|?*]'
    return re.sub(bad, "_", stem)


def main() -> int:
    if len(sys.argv) < 3:
        print("Usage: extract_to_md.py <out_dir> <file>...", file=sys.stderr)
        return 1
    out_dir = Path(sys.argv[1])
    out_dir.mkdir(parents=True, exist_ok=True)
    for src in map(Path, sys.argv[2:]):
        if not src.is_file():
            print(f"Skip (missing): {src}", file=sys.stderr)
            continue
        suf = src.suffix.lower()
        if suf == ".pptx":
            body = pptx_to_md(src)
        elif suf == ".pdf":
            body = pdf_to_md(src)
        else:
            print(f"Skip (unsupported): {src}", file=sys.stderr)
            continue
        out = out_dir / f"{sanitize_filename(src.stem)}.md"
        out.write_text(body, encoding="utf-8")
        print(f"Wrote {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
